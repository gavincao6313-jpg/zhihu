"""
extract_slides.py — 从已处理视频中提取幻灯片，同时输出 PDF 和 PPTX。

读取 cache/keyframes/<stem>/manifest.json 中的幻灯片切换事件，
对每个事件从原始 MP4 精确抽取全分辨率关键帧，去重后打包。

输出:
    Slides/<stem>/frames/   — 去重后的高清幻灯片 JPEG（用于检查）
    Slides/<stem>/slides.pdf
    Slides/<stem>/slides.pptx  （需要 python-pptx，否则跳过并提示）

用法:
    python extract_slides.py                          # 处理所有有缓存的视频
    python extract_slides.py --video "原理_A01.mp4"   # 指定单个视频
    python extract_slides.py --sim-threshold 5        # 调严去重（默认 8）
    python extract_slides.py --no-dedup               # 关闭去重
"""

import argparse
import json
import subprocess
from pathlib import Path

import numpy as np
from PIL import Image

# ── 路径配置 ─────────────────────────────────────────────────────────────────

VIDEOS_DIR         = Path(__file__).parent / "Videos"
CACHE_DIR          = Path(__file__).parent / "cache"
KEYFRAME_CACHE_DIR = CACHE_DIR / "keyframes"
SLIDES_DIR         = Path(__file__).parent / "Slides"
FFMPEG_TIMEOUT     = 30
VIDEO_EXTS         = {".mp4", ".mkv", ".avi", ".mov"}

# ── 帧提取 ───────────────────────────────────────────────────────────────────

def _extract_hires_frame(video_path: Path, timestamp_s: int, out_path: Path) -> bool:
    """从 video_path 的 timestamp_s 秒处抽取全分辨率 JPEG。"""
    try:
        subprocess.run([
            "ffmpeg",
            "-ss", str(timestamp_s),   # 输入端快速定位（不从头解码）
            "-i", str(video_path),
            "-vframes", "1",
            "-q:v", "2",               # JPEG 质量 2（接近无损）
            "-y", str(out_path),
        ], capture_output=True, check=True, timeout=FFMPEG_TIMEOUT)
        return out_path.exists() and out_path.stat().st_size > 0
    except Exception as exc:
        print(f"    [warn] 抽帧失败 ({timestamp_s}s): {exc}")
        return False


# ── 去重 ─────────────────────────────────────────────────────────────────────

def _thumb(path: Path) -> np.ndarray:
    return np.array(Image.open(path).resize((64, 36)).convert("L"), dtype=np.float32)


def _are_similar(a: Path, b: Path, threshold: float) -> bool:
    return float(np.mean(np.abs(_thumb(a) - _thumb(b)))) < threshold


def _deduplicate(frames: list[Path], threshold: float) -> list[Path]:
    """与已保留帧中任一相似则丢弃（O(n²)，n 通常 < 200）。"""
    kept: list[Path] = []
    for frame in frames:
        if not any(_are_similar(frame, prev, threshold) for prev in kept):
            kept.append(frame)
    return kept


# ── 事件加载 ─────────────────────────────────────────────────────────────────

def _load_slide_timestamps(stem: str) -> list[int]:
    """从 manifest.json 提取 slide 事件时间戳（秒），并在最前插入第 0 秒。"""
    manifest = KEYFRAME_CACHE_DIR / stem / "manifest.json"
    if not manifest.exists():
        return []
    with open(manifest, "r", encoding="utf-8") as f:
        data = json.load(f)
    slide_ts = [e["frame_idx"] for e in data.get("events", []) if e.get("type") == "slide"]
    return [0] + slide_ts


# ── 主流程 ───────────────────────────────────────────────────────────────────

def extract_slides(
    video_path: Path,
    sim_threshold: float = 8.0,
    dedup: bool = True,
) -> Path | None:
    stem = video_path.stem
    timestamps = _load_slide_timestamps(stem)

    if not timestamps:
        print(f"  [{stem}] 无关键帧缓存，跳过（先运行 zhihuTTS.py 处理该视频）")
        return None

    print(f"  [{stem}] 幻灯片事件: {len(timestamps)} 个（含初始帧）")

    out_dir    = SLIDES_DIR / stem
    frames_dir = out_dir / "frames"
    frames_dir.mkdir(parents=True, exist_ok=True)

    # ── Step 1: 高清抽帧 ──────────────────────────────────────────────────────
    extracted: list[Path] = []
    total = len(timestamps)
    for i, ts in enumerate(timestamps):
        out_frame = frames_dir / f"raw_{i+1:03d}_{ts}s.jpg"
        if _extract_hires_frame(video_path, ts, out_frame):
            extracted.append(out_frame)
        print(f"    抽帧进度: {i+1}/{total}", end="\r", flush=True)
    print(f"\n  [{stem}] 高清抽帧完成: {len(extracted)} / {total} 张")

    if not extracted:
        print(f"  [{stem}] 全部抽帧失败，退出")
        return None

    # ── Step 2: 去重 ─────────────────────────────────────────────────────────
    if dedup:
        deduped = _deduplicate(extracted, sim_threshold)
        removed = len(extracted) - len(deduped)
        if removed:
            print(f"  [{stem}] 去重: 移除 {removed} 张相似帧，保留 {len(deduped)} 张")
    else:
        deduped = extracted

    # 重命名为连续序号（方便人工检查）
    final_frames: list[Path] = []
    for i, src in enumerate(deduped):
        dst = frames_dir / f"slide_{i+1:03d}.jpg"
        src.rename(dst)
        final_frames.append(dst)

    # 清理未入选的 raw 帧
    for raw in frames_dir.glob("raw_*.jpg"):
        raw.unlink(missing_ok=True)

    # ── Step 3: 输出 PDF ──────────────────────────────────────────────────────
    pdf_path = out_dir / "slides.pdf"
    images = [Image.open(p).convert("RGB") for p in final_frames]
    images[0].save(
        pdf_path,
        save_all=True,
        append_images=images[1:],
        resolution=150,
    )
    print(f"  [{stem}] PDF  → {pdf_path}  ({len(final_frames)} 页)")

    # ── Step 4: 输出 PPTX ────────────────────────────────────────────────────
    pptx_path = out_dir / "slides.pptx"
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width  = Inches(13.33)   # 16:9 宽屏标准
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]       # 空白版式

        for img_path in final_frames:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(img_path),
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        prs.save(pptx_path)
        print(f"  [{stem}] PPTX → {pptx_path}  ({len(final_frames)} 页)")
    except ImportError:
        print("  [!] python-pptx 未安装，跳过 PPTX 输出")
        print("      安装命令: pip install python-pptx")

    return out_dir


# ── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="从已处理视频提取幻灯片，输出 PDF + PPTX",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "--video", metavar="FILE",
        help="指定视频文件名（含扩展名），位于 Videos/ 目录下",
    )
    parser.add_argument(
        "--sim-threshold", type=float, default=8.0, metavar="N",
        help="去重相似度阈值（0–255），越低越严格，默认 8.0",
    )
    parser.add_argument(
        "--no-dedup", action="store_true",
        help="关闭去重，输出全部切换帧（含老师翻回去的重复页）",
    )
    args = parser.parse_args()

    if args.video:
        video_path = VIDEOS_DIR / args.video
        if not video_path.exists():
            print(f"[错误] 找不到视频: {video_path}")
            return
        extract_slides(video_path,
                       sim_threshold=args.sim_threshold,
                       dedup=not args.no_dedup)
    else:
        candidates = sorted(
            p for p in VIDEOS_DIR.iterdir()
            if p.suffix.lower() in VIDEO_EXTS
            and (KEYFRAME_CACHE_DIR / p.stem / "manifest.json").exists()
        )
        if not candidates:
            print("没有找到带关键帧缓存的视频。请先运行 zhihuTTS.py 处理目标视频。")
            return
        print(f"找到 {len(candidates)} 个有缓存的视频\n")
        for video in candidates:
            extract_slides(video,
                           sim_threshold=args.sim_threshold,
                           dedup=not args.no_dedup)
            print()


if __name__ == "__main__":
    main()
